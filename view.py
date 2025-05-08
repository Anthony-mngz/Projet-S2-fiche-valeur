from pptx import Presentation
import comtypes.client
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import logging
import sys
import subprocess
import os
import yaml
import smtplib
from email.message import EmailMessage
from pathlib import Path

logger = logging.getLogger(__name__)

def update_ppt(template_path: str,
               output_path: str,
               data: pd.DataFrame,
               description: str,
               dictionnaire: dict,
               ratios: dict,
               chart_path: str,
               esg_df: pd.DataFrame,
               news_items: list[dict],
               calendar: dict,
               forecasts: dict):
    """Met à jour tous les slides du template PPTX."""
    # Remplace les clés par les valeurs
    prs = Presentation(template_path)
    logger.info("Chargement du template PPTX : %s", template_path)
    logger.info("Remplacement du placeholder share_name sur toutes les slides")

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for key, value in dictionnaire.items():
                    if key in shape.text:
                        shape.text = shape.text.replace(key, str(value))
                        if key == "shortName":
                            text_frame = shape.text_frame
                            for paragraph in text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.CENTER  # Centre le texte
                                for run in paragraph.runs:
                                    run.font.color.rgb = RGBColor(255, 255, 255)  # Change la couleur en blanc


    logger.info("Remplissage de la slide 1 (Infos & Financials)")
    # Slide 1: Infos & Financials
    slide = prs.slides[0]

    for shape in slide.shapes:
        if shape.has_text_frame:
            for key, value in ratios.items():
                if value is None:
                    value = "N/A"
                if key in shape.text:
                    shape.text = shape.text.replace(key, str(value))

    for shape in slide.shapes:
        if hasattr(shape, "text") and "description" in shape.text:
            shape.text = description
            if hasattr(shape, "text_frame") and shape.text_frame:
                shape.text_frame.text = description
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.JUSTIFY
                    for run in paragraph.runs:
                        run.font.size = Pt(9)

    if chart_path:
        slide.shapes.add_picture(chart_path, Inches(-0.4), Inches(0.7), width=Inches(7.7),
                                 height=Inches(3.9))

    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table

            # Safe conversion for "Total Revenue"
            try:
                revenue_n1 = int(data.loc["Total Revenue", "N"] / 1000) if pd.notna(
                    data.loc["Total Revenue", "N"]) else "N/A"
            except KeyError:
                revenue_n1 = "N/A"

            try:
                revenue_n2 = int(data.loc["Total Revenue", "N-1"] / 1000) if pd.notna(
                    data.loc["Total Revenue", "N-1"]) else "N/A"
            except KeyError:
                revenue_n2 = "N/A"

            try:
                revenue_n3 = int(data.loc["Total Revenue", "N-2"] / 1000) if pd.notna(
                    data.loc["Total Revenue", "N-2"]) else "N/A"
            except KeyError:
                revenue_n3 = "N/A"

            table.cell(1, 1).text = f"{revenue_n1:,}" if revenue_n1 != "N/A" else "N/A"
            table.cell(1, 2).text = f"{revenue_n2:,}" if revenue_n2 != "N/A" else "N/A"
            table.cell(1, 3).text = f"{revenue_n3:,}" if revenue_n3 != "N/A" else "N/A"

            # Safe conversion for "EBITDA"
            try:
                ebitda_n1 = int(data.loc["EBITDA", "N-1"] / 1000) if pd.notna(data.loc["EBITDA", "N-1"]) else "N/A"
            except KeyError:
                ebitda_n1 = "N/A"

            try:
                ebitda_n2 = int(data.loc["EBITDA", "N-2"] / 1000) if pd.notna(data.loc["EBITDA", "N-2"]) else "N/A"
            except KeyError:
                ebitda_n2 = "N/A"

            try:
                ebitda_n3 = int(data.loc["EBITDA", "N-3"] / 1000) if pd.notna(data.loc["EBITDA", "N-3"]) else "N/A"
            except KeyError:
                ebitda_n3 = "N/A"

            table.cell(2, 1).text = f"{ebitda_n1:,}" if ebitda_n1 != "N/A" else "N/A"
            table.cell(2, 2).text = f"{ebitda_n2:,}" if ebitda_n2 != "N/A" else "N/A"
            table.cell(2, 3).text = f"{ebitda_n3:,}" if ebitda_n3 != "N/A" else "N/A"

            # Safe conversion for "EBIT"
            try:
                ebit_n1 = int(data.loc["EBIT", "N-1"] / 1000) if pd.notna(data.loc["EBIT", "N-1"]) else "N/A"
            except KeyError:
                ebit_n1 = "N/A"

            try:
                ebit_n2 = int(data.loc["EBIT", "N-2"] / 1000) if pd.notna(data.loc["EBIT", "N-2"]) else "N/A"
            except KeyError:
                ebit_n2 = "N/A"

            try:
                ebit_n3 = int(data.loc["EBIT", "N-3"] / 1000) if pd.notna(data.loc["EBIT", "N-3"]) else "N/A"
            except KeyError:
                ebit_n3 = "N/A"

            table.cell(3, 1).text = f"{ebit_n1:,}" if ebit_n1 != "N/A" else "N/A"
            table.cell(3, 2).text = f"{ebit_n2:,}" if ebit_n2 != "N/A" else "N/A"
            table.cell(3, 3).text = f"{ebit_n3:,}" if ebit_n3 != "N/A" else "N/A"

            # Safe conversion for "Net Debt on EBITDA"
            try:
                debt_on_ebitda_n1 = round(float((data.loc["Net Debt", "N-1"] / data.loc["EBITDA", "N-1"])),
                                          2) if pd.notna(
                    data.loc["EBITDA", "N-1"] or data.loc["Net Debt", "N-1"]) else "N/A"
            except KeyError:
                debt_on_ebitda_n1 = "N/A"

            try:
                debt_on_ebitda_n2 = round(float((data.loc["Net Debt", "N-2"] / data.loc["EBITDA", "N-2"])),
                                          2) if pd.notna(
                    data.loc["EBITDA", "N-2"] or data.loc["Net Debt", "N-2"]) else "N/A"
            except KeyError:
                debt_on_ebitda_n2 = "N/A"

            try:
                debt_on_ebitda_n3 = round(float((data.loc["Net Debt", "N-3"] / data.loc["EBITDA", "N-3"])),
                                          2) if pd.notna(
                    data.loc["EBITDA", "N-3"] or data.loc["Net Debt", "N-3"]) else "N/A"
            except KeyError:
                debt_on_ebitda_n3 = "N/A"

            table.cell(4, 1).text = f"{debt_on_ebitda_n1:,}x" if debt_on_ebitda_n1 != "N/A" else "N/A"
            table.cell(4, 2).text = f"{debt_on_ebitda_n2:,}x" if debt_on_ebitda_n2 != "N/A" else "N/A"
            table.cell(4, 3).text = f"{debt_on_ebitda_n3:,}x" if debt_on_ebitda_n3 != "N/A" else "N/A"

            # Safe conversion for "EPS"
            try:
                diluted_eps_n1 = float(data.loc["Diluted EPS", "N-1"]) if pd.notna(
                    data.loc["Diluted EPS", "N-1"]) else "N/A"
            except KeyError:
                diluted_eps_n1 = "N/A"

            try:
                diluted_eps_n2 = float(data.loc["Diluted EPS", "N-2"]) if pd.notna(
                    data.loc["Diluted EPS", "N-2"]) else "N/A"
            except KeyError:
                diluted_eps_n2 = "N/A"

            try:
                diluted_eps_n3 = round(float(data.loc["Diluted EPS", "N-3"]), 2) if pd.notna(
                    data.loc["Diluted EPS", "N-3"]) else "N/A"
            except KeyError:
                diluted_eps_n3 = "N/A"

            table.cell(5, 1).text = f"{diluted_eps_n1:,}" if diluted_eps_n1 != "N/A" else "N/A"
            table.cell(5, 2).text = f"{diluted_eps_n2:,}" if diluted_eps_n2 != "N/A" else "N/A"
            table.cell(5, 3).text = f"{diluted_eps_n3:,}" if diluted_eps_n3 != "N/A" else "N/A"

    logger.info("Remplissage de la slide 2 (ESG + News)")
    # Slide 2: ESG + News
    slide = prs.slides[1]
    # Formattage ESG
    def fmt_esg(index_name: str) -> str:
        try:
            raw = esg_df.loc[index_name, "esgScores"]
        except (KeyError, AttributeError):
            return "N/A"
        if isinstance(raw, (list, tuple)):
            sep = "\n" if index_name == "relatedControversy" else " / "   #Cela devait servir à séparer les controverses mais j'ai essayé plusieurs méthodes sans y arriver
            return sep.join(str(x) for x in raw)
        return str(raw)

    total_esg   = fmt_esg("totalEsg")
    month       = fmt_esg("ratingMonth")
    year        = fmt_esg("ratingYear")
    controversy = fmt_esg("highestControversy")
    related     = fmt_esg("relatedControversy")

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text.strip().lower()
        if "total-esg" in txt:
            shape.text = total_esg
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(40)
        if "month" in txt:
            shape.text = month
        if "year" in txt:
            shape.text = year
        if "level_contro" in txt:
            shape.text = controversy
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(40)
        if "related_contro" in txt:
            shape.text = related

    # Formatage News
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        if tf.text.strip().lower() == "news":
            tf.clear()
            for item in news_items[:9]:
                para = tf.add_paragraph()
                # Récupération du titre dans content
                content = item.get("content", {})
                title = content.get("title", "")
                # Date préférentiellement dans 'datetime'
                dt = item.get("datetime")
                if not dt and "providerPublishTime" in item:
                    try:
                        dt = datetime.fromtimestamp(item.get("providerPublishTime"))
                    except Exception:
                        dt = None
                date_str = dt.strftime('%Y-%m-%d') if hasattr(dt, 'strftime') else ''
                para.text = f"{date_str} - {title}"
                para.level = 0


    logger.info("Remplissage de la slide 3 (Consensus, Calendar, Upside/Downside)")
    # Slide 3: Economic Calendar, Price Consensus, Upside/Downside
    slide = prs.slides[2]
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            # Price Consensus
            if 'price_consensus' in tf.text.lower():
                mn = forecasts.get('target_mean_price', 'N/A')
                lo = forecasts.get('target_low_price', 'N/A')
                hi = forecasts.get('target_high_price', 'N/A')
                rk = forecasts.get('recommendation_key', '')
                tf.text = f"low {lo}, mean {mn}, high {hi}"
            if 'reco_key' in tf.text.lower():
                rk = forecasts.get('recommendation_key', '')
                tf.text = rk
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(40)
    # Calendar
    if isinstance(calendar, dict):
        events = list(calendar.items())
        for idx, (name, date) in enumerate(events[:3], start=1):
            # si date est une liste, on en prend le premier élément
            if isinstance(date, (list, tuple)) and len(date) > 0:
                d0 = date[0]
            else:
                d0 = date
            # formate la date si possible
            if hasattr(d0, "strftime"):
                date_str = d0.strftime("%Y-%m-%d")
            else:
                date_str = str(d0)

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                txt = shape.text.strip().lower()
                if f"next_event{idx}" in txt:
                    shape.text = name
                if f"date_next_event{idx}" in txt:
                    shape.text = date_str

    # Upside/Downside
    cp = dictionnaire.get('previousClose') or None
    if not cp and 'N' in data.columns:
        cp = data['N'].iloc[-1]
    up = None
    if cp and forecasts.get('target_mean_price'):
        up = round((forecasts['target_mean_price']/cp - 1)*100, 2)
    for shape in slide.shapes:
        if shape.has_text_frame and 'upside_downside' in shape.text.lower():
            shape.text = f"{up}%" if up is not None else "N/A"
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(40)

    logger.info("Enregistrement du PPT mis à jour dans %s", output_path)
    prs.save(output_path)

def convert_ppt_to_pdf(ppt_path: str, pdf_path: str):
    """Convertit un fichier PowerPoint en PDF via COM."""
    logger.info("Conversion PPTX en PDF : %s → %s", ppt_path, pdf_path)
    try :
        import comtypes.client
        try:
            ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
            ppt_app.Visible = True
            presentation = ppt_app.Presentations.Open(ppt_path)
            presentation.SaveAs(pdf_path, 32)  # 32 = PDF
            presentation.Close()
            ppt_app.Quit()
        except Exception as e:
            raise RuntimeError(f"Erreur conversion PPT en PDF: {e}")
        logger.info("PDF créé avec succès")
    except Exception as e:
        logger.error("Erreur conversion PPT→PDF: %s", e)
        raise

def display_pdf(pdf_path: str):
    """
    Ouvre le PDF généré avec le lecteur par défaut du système.
    - Windows : os.startfile
    - macOS   : open
    - Linux   : xdg-open
    """
    try:
        logger.info(f"Ouverture du PDF : {pdf_path}")
        if sys.platform.startswith("win"):
            os.startfile(pdf_path)
        elif sys.platform == "darwin":
            subprocess.run(["open", pdf_path], check=False)
        else:
            subprocess.run(["xdg-open", pdf_path], check=False)
    except Exception as e:
        logger.error(f"Impossible d'ouvrir le PDF : {e}")

def send_report_via_email(pdf_path: str, cfg: dict, ticker):
    """
    Envoie le PDF en pièce jointe via SMTP (smtplib).
    Lit les paramètres (serveur, port, user/pwd, destinataire, sujet, corps) depuis le dict cfg.
    """
    # Prépare le message
    subject   = cfg["email_subject"].format(ticker=ticker)
    body      = cfg["email_body"].format(ticker=ticker)
    recipient = cfg["email_recipient"]

    msg = EmailMessage()
    msg["Subject"] = subject
    msg["From"]    = cfg["smtp_user"]
    msg["To"]      = recipient
    msg.set_content(body)

    # Ajoute le PDF en pièce jointe
    pdf_file = Path(pdf_path)
    if not pdf_file.exists():
        raise FileNotFoundError(f"Le fichier PDF n'existe pas : {pdf_path}")
    with pdf_file.open("rb") as f:
        data = f.read()
    msg.add_attachment(
        data,
        maintype="application",
        subtype="pdf",
        filename=pdf_file.name
    )

    # Envoi via SMTP_SSL
    server = cfg["smtp_server"]
    port   = cfg["smtp_port"]
    user   = cfg["smtp_user"]
    pwd    = cfg["smtp_password"]

    logger.info(f"Envoi de l'email à {recipient} via {server}:{port}")
    with smtplib.SMTP_SSL(server, port) as smtp:
        smtp.login(user, pwd)
        smtp.send_message(msg)
    logger.info("Email envoyé avec succès.")